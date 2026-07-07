"""auto_regen.py — _trigger.json 감지 → NLM 재생성 → 이미지 반영 자동 루프

브라우저(image_selector.py)에서 "Claude에게 재생성 요청" 클릭 시:
  1. _trigger.json 생성됨
  2. 이 스크립트가 감지
  3. NLM에서 새 슬라이드 생성 + PPTX 다운로드
  4. 이미지 추출 → 해당 슬롯에 새 세트로 저장
  5. _trigger.json 삭제 → 브라우저 자동 새로고침

사용: python scripts/auto_regen.py [--notebook <notebook-id>] [--dir <설교폴더>]
"""

import sys
import os
import io
import json
import time
import re
import subprocess
import shutil
from pathlib import Path
from datetime import datetime

# Windows cp949 인코딩 이슈 방지
if sys.platform == "win32":
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")
    sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding="utf-8", errors="replace")

# ──────────────────────────────────────────────
# 설정
# ──────────────────────────────────────────────
POLL_INTERVAL = 5  # 초
NLM_ENV = {
    "PYTHONIOENCODING": "utf-8",
    "PYTHONLEGACYWINDOWSSTDIO": "utf-8",
}

# PPT 슬롯별 이미지 설명 (NLM focus 프롬프트용)
# 시편 95:1-11 "마음의 예배" — 2026.04.15 수요기도회
SLOT_FOCUS = {
    "02": "Close-up of a singer holding a microphone on a brightly lit concert stage, dramatic spotlight, mouth movement frozen mid-song, cinematic film noir lighting, the contrast between performance and authenticity",
    "03": "Two contrasting portraits side by side: one performer with empty hollow eyes singing on stage, the other with passionate weeping eyes singing live, dramatic stage lighting, cinematic emotional contrast",
    "08": "Vast cosmic creation scene, deep ocean trenches and towering mountain peaks, God's hands forming earth and sea, dramatic golden divine light piercing through dark clouds, epic scale, movie concept art",
    "10": "A shepherd in ancient Middle Eastern robes leading a flock of sheep through a golden sunlit field, the sheep gathered around him trustingly, intimate pastoral scene, warm golden hour lighting, cinematic photorealism",
    "14": "Wilderness wandering scene of the Israelites in the desert at Meribah, hardened faces looking away from Moses, dust and heat, dramatic biblical epic lighting, a sense of stubborn refusal and lost direction",
    "19": "Jesus the Good Shepherd standing among his sheep, the sheep recognizing his voice and turning toward him, golden divine light surrounding the scene, dark dramatic background with warm highlights, cinematic biblical art",
    "23": "Silhouette of Jesus Christ standing at a wooden door, hand raised gently knocking, warm golden light spilling from within, mysterious darkness outside, cinematic dramatic lighting, emotional and intimate atmosphere",
    "24": "Close-up of a hand gently knocking on an old wooden door from outside, warm golden light spilling through cracks in the door, mysterious darkness behind the hand, cinematic lighting reminiscent of Holman Hunt's Light of the World",
}

NLM_BASE_FOCUS = (
    "ABSOLUTELY NO TEXT anywhere. NO Korean letters. NO English letters. NO Greek letters. NO Hebrew letters. "
    "ZERO words on any slide. NO captions, NO titles, NO labels, NO subtitles. "
    "CINEMATIC PHOTOREALISM ONLY. Dark dramatic lighting with golden warm highlights. "
    "ADULT figures only, no children. Movie concept art quality. "
    "ALL HUMAN FIGURES MUST BE KOREAN PEOPLE (East Asian Korean ethnicity). "
    "Korean facial features, Korean appearance, Korean skin tone. "
    "Even in biblical or worship contexts, depict figures as modern Korean believers, "
    "Korean worshippers, Korean Christians praying, listening, or kneeling. "
    "NO Western people, NO non-Asian figures unless a specific historical biblical character requires it. "
    "When in doubt, make them Korean. "
    "Each slide must be a SINGLE powerful photorealistic scene."
)


# ──────────────────────────────────────────────
# 경로 탐지
# ──────────────────────────────────────────────
def find_sermon_dir():
    """설교 폴더 자동 탐지 (--dir 인자 또는 가장 최근 폴더)"""
    for i, arg in enumerate(sys.argv):
        if arg == "--dir" and i + 1 < len(sys.argv):
            return Path(sys.argv[i + 1]).resolve()
    project = Path(__file__).parent.parent.resolve()
    dirs = sorted(
        [d for d in project.iterdir() if d.is_dir() and d.name[:8].isdigit()],
        reverse=True,
    )
    return dirs[0] if dirs else None


def find_notebook_id():
    """--notebook 인자 또는 설교폴더/.nlm_notebook 파일에서 노트북 ID 읽기"""
    for i, arg in enumerate(sys.argv):
        if arg == "--notebook" and i + 1 < len(sys.argv):
            return sys.argv[i + 1]
    # 설교폴더에 저장된 ID
    nb_file = SERMON_DIR / ".nlm_notebook"
    if nb_file.exists():
        return nb_file.read_text(encoding="utf-8").strip()
    return None


# ──────────────────────────────────────────────
# 세트 프리픽스 관리
# ──────────────────────────────────────────────
def get_existing_prefixes(images_dir):
    """images/ 폴더에서 기존 세트 프리픽스 목록"""
    prefixes = set()
    for f in images_dir.glob("*.png"):
        if f.name.startswith("ppt_") or f.name.startswith("_"):
            continue
        parts = f.stem.rsplit("_", 1)
        if len(parts) == 2 and parts[1].isdigit():
            prefixes.add(parts[0])
    return sorted(prefixes)


def get_next_prefix(images_dir):
    """다음 사용할 세트 프리픽스 (A→B→...→Z→AA→AB...)"""
    existing = set(get_existing_prefixes(images_dir))
    # 단일 문자 먼저
    for c in "ABCDEFGHIJKLMNOPQRSTUVWXYZ":
        if c not in existing:
            return c
    # 두 글자
    for c1 in "ABCDEFGHIJKLMNOPQRSTUVWXYZ":
        for c2 in "ABCDEFGHIJKLMNOPQRSTUVWXYZ":
            combo = c1 + c2
            if combo not in existing:
                return combo
    return f"S{len(existing)}"


# ──────────────────────────────────────────────
# NLM 슬라이드 생성 + 다운로드
# ──────────────────────────────────────────────
def run_nlm(args, timeout=180):
    """NLM CLI 실행 (인코딩 처리 포함)"""
    env = {**os.environ, **NLM_ENV}
    cmd = ["nlm"] + args
    log(f"  $ {' '.join(cmd[:6])}...")
    try:
        result = subprocess.run(
            cmd,
            capture_output=True,
            text=True,
            encoding="utf-8",
            errors="replace",
            env=env,
            timeout=timeout,
        )
        if result.stdout and result.stdout.strip():
            for line in result.stdout.strip().split("\n")[:5]:
                log(f"    {line}")
        if result.returncode != 0 and result.stderr:
            log(f"  [ERROR] {result.stderr.strip()[:300]}")
        return result
    except subprocess.TimeoutExpired:
        log(f"  [TIMEOUT] {timeout}s exceeded")
        return None
    except Exception as e:
        log(f"  [EXCEPTION] {e}")
        return None


def parse_artifact_id(output_text):
    """NLM slides create 출력에서 Artifact ID 추출"""
    match = re.search(r"Artifact ID:\s*([a-f0-9\-]+)", output_text)
    return match.group(1) if match else None


def wait_for_completion(notebook_id, artifact_id, max_wait=600, poll_interval=15):
    """NLM 슬라이드 생성 완료까지 대기 (비동기 폴링)

    Args:
        max_wait: 최대 대기 시간 (초, 기본 5분)
        poll_interval: 폴링 간격 (초, 기본 10초)
    Returns:
        True if completed, False if failed/timeout
    """
    elapsed = 0
    while elapsed < max_wait:
        r = run_nlm(["studio", "status", notebook_id], timeout=30)
        if r is None:
            time.sleep(poll_interval)
            elapsed += poll_interval
            continue

        try:
            artifacts = json.loads(r.stdout)
            for art in artifacts:
                if art.get("id") == artifact_id:
                    status = art.get("status", "")
                    if status == "completed":
                        log(f"  [OK] Slides completed! ({elapsed}s)")
                        return True
                    elif status in ("failed", "error"):
                        log(f"  [FAIL] Slide generation failed: {status}")
                        return False
                    else:
                        log(f"  [WAIT] Generating... ({elapsed}s, status: {status})")
                    break
        except (json.JSONDecodeError, TypeError):
            pass

        time.sleep(poll_interval)
        elapsed += poll_interval

    log(f"  [TIMEOUT] Exceeded {max_wait}s")
    return False


def generate_and_download(notebook_id, slots, output_pptx, user_prompts=None):
    """NLM 슬라이드 생성(비동기) → 완료 대기 → PPTX 다운로드"""
    user_prompts = user_prompts or {}

    # 슬롯별 포커스 프롬프트 조합 (기본 + 사용자 커스텀)
    scene_parts = []
    for s in slots:
        base = SLOT_FOCUS.get(s, "")
        user = user_prompts.get(s, "")
        if user:
            scene_parts.append(f"Scene for PPT{s}: {base}. USER REQUEST: {user}")
            log(f"  Slot {s} custom prompt: {user}")
        elif base:
            scene_parts.append(f"Scene for PPT{s}: {base}")
    focus = NLM_BASE_FOCUS
    if scene_parts:
        focus += " " + "; ".join(scene_parts)

    # 1) 슬라이드 생성 시작 (비동기)
    log(f"  NLM slides creating... (slots: {', '.join(slots)})")
    r = run_nlm([
        "slides", "create", notebook_id,
        "--language", "ko",
        "--focus", focus,
        "--confirm",
    ], timeout=180)
    if r is None or r.returncode != 0:
        return False

    # 2) Artifact ID 추출
    artifact_id = parse_artifact_id(r.stdout)
    if not artifact_id:
        log(f"  [WARN] Artifact ID not found. Waiting 60s then downloading...")
        time.sleep(60)
    else:
        log(f"  Artifact ID: {artifact_id}")
        # 3) 완료까지 폴링 대기
        if not wait_for_completion(notebook_id, artifact_id):
            log(f"  [FAIL] Completion wait failed")
            return False

    # 4) PPTX 다운로드 (--id로 특정 아티팩트 지정)
    log(f"  Downloading PPTX...")
    dl_args = [
        "download", "slide-deck", notebook_id,
        "--format", "pptx",
        "--output", str(output_pptx),
        "--no-progress",
    ]
    if artifact_id:
        dl_args += ["--id", artifact_id]
    r = run_nlm(dl_args, timeout=120)
    if r is None or r.returncode != 0:
        return False

    return output_pptx.exists()


# ──────────────────────────────────────────────
# PPTX → 이미지 추출
# ──────────────────────────────────────────────
def extract_images_from_pptx(pptx_path):
    """PPTX에서 슬라이드 이미지 추출 (큰 이미지만)"""
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        log("  [ERROR] python-pptx not installed! pip install python-pptx")
        return []

    prs = Presentation(str(pptx_path))
    images = []

    for slide_idx, slide in enumerate(prs.slides):
        slide_images = []
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                try:
                    blob = shape.image.blob
                    ct = shape.image.content_type
                    if len(blob) > 30000:  # 30KB 이상만 (아이콘/로고 제외)
                        slide_images.append({
                            "slide_idx": slide_idx,
                            "blob": blob,
                            "content_type": ct,
                            "size": len(blob),
                        })
                except Exception:
                    pass

        # 슬라이드당 가장 큰 이미지만 선택
        if slide_images:
            biggest = max(slide_images, key=lambda x: x["size"])
            images.append(biggest)

    log(f"  Extracted: {len(images)} images (total slides: {len(prs.slides)})")
    return images


def map_images_to_slots(images, slots, total_ppt_slots=14):
    """추출된 이미지를 요청된 슬롯에 위치 기반으로 매핑

    매핑 전략: 슬라이드 위치(0~1) ↔ PPT 슬롯 위치(0~1)를 비교하여
    가장 가까운 이미지를 각 슬롯에 배정
    """
    if not images or not slots:
        return {}

    total_imgs = len(images)
    mapping = {}
    used = set()

    # 각 슬롯에 대해 가장 가까운 이미지 찾기
    for slot in sorted(slots, key=lambda s: int(s)):
        slot_pos = (int(slot) - 1) / max(total_ppt_slots - 1, 1)

        best_idx = None
        best_dist = float("inf")
        for i, img in enumerate(images):
            if i in used:
                continue
            img_pos = img["slide_idx"] / max(total_imgs - 1, 1) if total_imgs > 1 else 0.5
            dist = abs(img_pos - slot_pos)
            if dist < best_dist:
                best_dist = dist
                best_idx = i

        if best_idx is not None:
            mapping[slot] = images[best_idx]
            used.add(best_idx)

    return mapping


def save_mapped_images(mapping, images_dir, prefix):
    """매핑된 이미지를 {prefix}_{slot}.png로 저장"""
    saved = []
    for slot, img_data in mapping.items():
        ext = "png"
        ct = img_data.get("content_type", "")
        if "jpeg" in ct or "jpg" in ct:
            ext = "jpg"

        out_path = images_dir / f"{prefix}_{slot}.{ext}"
        # PNG가 아니면 변환
        if ext != "png":
            try:
                from PIL import Image
                import io
                pil_img = Image.open(io.BytesIO(img_data["blob"]))
                pil_img.save(str(out_path.with_suffix(".png")), "PNG")
                out_path = out_path.with_suffix(".png")
            except ImportError:
                out_path.write_bytes(img_data["blob"])
        else:
            out_path.write_bytes(img_data["blob"])

        size_kb = len(img_data["blob"]) // 1024
        log(f"  [OK] PPT {slot} -> {out_path.name} ({size_kb}KB)")
        saved.append(slot)

    return saved


# ──────────────────────────────────────────────
# 트리거 처리
# ──────────────────────────────────────────────
def process_trigger(trigger_file, images_dir, notebook_id):
    """_trigger.json 감지 → NLM 생성 + 추출 → _review.json 작성 (매핑은 Claude가 수행)"""
    trigger = json.loads(trigger_file.read_text(encoding="utf-8"))
    slots = trigger.get("slots", [])
    if not slots:
        log("  No slots, removing trigger")
        trigger_file.unlink()
        return

    log(f"\n{'='*55}")
    log(f"[REGEN] Slots: PPT {', '.join(slots)}")
    log(f"{'='*55}")

    # 상태 → processing (브라우저에서 대기 화면 표시)
    trigger["status"] = "processing"
    trigger_file.write_text(
        json.dumps(trigger, ensure_ascii=False, indent=2), encoding="utf-8"
    )

    # 새 세트 프리픽스
    prefix = get_next_prefix(images_dir)
    log(f"  New set: {prefix}")

    # 사용자 프롬프트 읽기
    user_prompts = trigger.get("prompts", {})

    # NLM 생성 + 다운로드
    pptx_path = images_dir / f"nlm_{prefix}.pptx"
    success = generate_and_download(notebook_id, slots, pptx_path, user_prompts)

    if not success:
        log("  [FAIL] NLM generation failed! Removing trigger.")
        cleanup_trigger(trigger_file)
        return

    # 이미지 추출 (슬롯 매핑 없이 순번으로 저장)
    images = extract_images_from_pptx(pptx_path)
    if not images:
        log("  [FAIL] Image extraction failed! Removing trigger.")
        cleanup_trigger(trigger_file)
        return

    # _pending/ 폴더에 순번으로 저장
    pending_dir = images_dir / "_pending"
    pending_dir.mkdir(exist_ok=True)
    raw_files = []
    for i, img_data in enumerate(images):
        fname = f"{prefix}_raw_{i:02d}.png"
        out_path = pending_dir / fname
        out_path.write_bytes(img_data["blob"])
        raw_files.append(fname)
        log(f"  raw {i:02d}: {fname} ({img_data['size']//1024}KB)")

    # _review.json 작성 → Claude가 이 파일을 감지하여 시각 검수 + 슬롯 매핑
    review = {
        "prefix": prefix,
        "slots": slots,
        "prompts": user_prompts,
        "raw_images": raw_files,
        "raw_dir": "_pending",
        "status": "needs_review",
        "timestamp": datetime.now().isoformat(),
    }
    review_file = images_dir / "_review.json"
    review_file.write_text(json.dumps(review, ensure_ascii=False, indent=2), encoding="utf-8")

    # 트리거 상태를 review로 변경 (브라우저에 "Claude 검수 중" 표시)
    trigger["status"] = "reviewing"
    trigger_file.write_text(
        json.dumps(trigger, ensure_ascii=False, indent=2), encoding="utf-8"
    )

    log(f"\n  [READY] {len(raw_files)} raw images extracted -> _pending/")
    log(f"  [WAIT] _review.json written. Waiting for Claude to map slots...")
    log(f"  Claude will read images, compare with manuscript, and assign to correct slots.\n")

    # Claude가 매핑 완료 후 _review.json과 _trigger.json을 삭제함
    # auto_regen은 여기서 대기 상태로 복귀


def cleanup_trigger(trigger_file):
    """트리거 파일 삭제"""
    if trigger_file.exists():
        trigger_file.unlink()


# ──────────────────────────────────────────────
# 로깅
# ──────────────────────────────────────────────
def log(msg):
    ts = datetime.now().strftime("%H:%M:%S")
    print(f"[{ts}] {msg}", flush=True)


# ──────────────────────────────────────────────
# 메인 루프
# ──────────────────────────────────────────────
def main():
    global SERMON_DIR

    SERMON_DIR = find_sermon_dir()
    if not SERMON_DIR or not SERMON_DIR.exists():
        print("[ERROR] Sermon dir not found. Use --dir option.")
        sys.exit(1)

    images_dir = SERMON_DIR / "images"
    if not images_dir.exists():
        images_dir.mkdir(parents=True)

    trigger_file = images_dir / "_trigger.json"
    notebook_id = find_notebook_id()

    if not notebook_id:
        print("[ERROR] NLM notebook ID required.")
        print("  Option 1: --notebook <ID>")
        print(f"  Option 2: save ID to {SERMON_DIR / '.nlm_notebook'}")
        sys.exit(1)

    # 기존 세트 표시
    existing = get_existing_prefixes(images_dir)

    print(f"""
===================================================
  [auto_regen] Image Auto-Regeneration Watcher
===================================================
  Sermon : {SERMON_DIR.name}
  Notebook: {notebook_id[:36]}
  Sets   : {', '.join(existing) if existing else 'none'}
  Poll   : {POLL_INTERVAL}s
===================================================

  Browser "regen request" -> auto NLM -> refresh
  Ctrl+C to stop
""")

    # 시작 시 이미 트리거가 있으면 바로 처리
    if trigger_file.exists():
        try:
            t = json.loads(trigger_file.read_text(encoding="utf-8"))
            if t.get("status") == "requested":
                process_trigger(trigger_file, images_dir, notebook_id)
        except Exception as e:
            log(f"Existing trigger error: {e}")
            cleanup_trigger(trigger_file)

    # 폴링 루프
    log("Waiting for regen requests...")
    while True:
        try:
            if trigger_file.exists():
                trigger = json.loads(trigger_file.read_text(encoding="utf-8"))
                if trigger.get("status") == "requested":
                    process_trigger(trigger_file, images_dir, notebook_id)
                    log("Waiting for next regen request...")
            time.sleep(POLL_INTERVAL)
        except KeyboardInterrupt:
            log("\nStopping.")
            break
        except Exception as e:
            log(f"Error: {e}")
            cleanup_trigger(trigger_file)
            time.sleep(POLL_INTERVAL)


if __name__ == "__main__":
    main()
