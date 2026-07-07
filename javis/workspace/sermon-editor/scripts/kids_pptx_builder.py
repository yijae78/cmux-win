"""아동부 PPT 슬라이드 빌더 — 19개 실작업 PPT 분석 기반 (장년부와 완전 독립)

슬라이드 유형 10종:
1. cover_scripture   — 표지: "성경 말씀" + 책명 + 장절
2. cover_text        — 성경 본문 전문 표시
3. cover_title       — "말씀 제목" + 설교 제목
4. full_image        — 풀스크린 이미지 (NLM)
5. text_image        — 풀스크린 텍스트 PNG (generate_text_slides.py)
6. question          — 거대 ? 효과
7. original_word     — 원어 제목 슬라이드 (헬라어/히브리어)
8. progressive_text  — 점진 노출 (같은 본문, 다른 강조)
9. image_with_masks  — 이미지 + 빨간 마스크 리빌
10. blank            — 빈 슬라이드
"""

import os
from lxml import etree
from pptx.util import Emu, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from kids_pptx_design import (
    SLIDE_WIDTH, SLIDE_HEIGHT,
    FONT_BODY, FONT_HEADER, FONT_GREEK, FONT_QUESTION,
    WHITE, YELLOW, CYAN, RED, BLACK,
    SIZE_BOOK_NAME, SIZE_TITLE, SIZE_CHAPTER_VERSE, SIZE_BODY,
    SIZE_BODY_SUB, SIZE_HEADER_TAG, SIZE_REF_TAG,
    SIZE_QUESTION_L, SIZE_QUESTION_XL,
    HEADER_BAR_XML,
    COVER_TAG_POS, COVER_BOOK_POS, SCRIPTURE_BOX_POS,
    SCRIPTURE_TAG_POS, TITLE_POS, SERIES_TAG_POS,
    QUESTION_POS, QUESTION_HEIGHT_MAP, IMAGE_FULLSCREEN,
    LAYOUT_BLANK,
    CHAR_SPACING_BODY, CHAR_SPACING_HEADER, CHAR_SPACING_QUESTION,
    LINE_SPACING_NORMAL, LINE_SPACING_WIDE,
    SPACE_BEFORE, SPACE_AFTER,
    SCRIPTURE_FRAME_POS,
)
from kids_pptx_animations import (
    inject_fade_in,
    inject_appear_colorpulse,
    inject_question_animation,
    inject_multi_mask_animation,
    inject_multi_fade_in,
    inject_click_fade_in,
    inject_text_mask_reveal,
)

NS_P = '{http://schemas.openxmlformats.org/presentationml/2006/main}'


# ============================================================
# 유틸리티
# ============================================================

def _new_slide(prs):
    """빈 화면 레이아웃으로 새 슬라이드 생성"""
    layouts = prs.slide_layouts
    # LAYOUT_BLANK 인덱스 시도, 실패하면 마지막 레이아웃
    try:
        layout = layouts[LAYOUT_BLANK]
    except IndexError:
        layout = layouts[len(layouts) - 1]

    slide = prs.slides.add_slide(layout)

    # 기본 플레이스홀더 모두 제거
    for ph in list(slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)

    return slide


def _apply_char_spacing(run, font_name):
    """폰트별 자간(character spacing) 적용 — OXML rPr spc 속성"""
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    rPr = run._r.find(f'{{{ns_a}}}rPr')
    if rPr is None:
        rPr = etree.SubElement(run._r, f'{{{ns_a}}}rPr')

    if font_name == FONT_BODY:
        rPr.set('spc', str(CHAR_SPACING_BODY))
    elif font_name == FONT_HEADER:
        rPr.set('spc', str(CHAR_SPACING_HEADER))
    elif font_name == FONT_QUESTION:
        rPr.set('spc', str(CHAR_SPACING_QUESTION))


def _apply_line_spacing(paragraph, line_spacing_pct=None, space_before=None, space_after=None):
    """단락 행간·전후 간격 적용 — OXML pPr"""
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    pPr = paragraph._p.find(f'{{{ns_a}}}pPr')
    if pPr is None:
        pPr = etree.SubElement(paragraph._p, f'{{{ns_a}}}pPr')
        # pPr을 첫 번째 자식으로 이동
        paragraph._p.remove(pPr)
        paragraph._p.insert(0, pPr)

    if line_spacing_pct is not None:
        lnSpc = etree.SubElement(pPr, f'{{{ns_a}}}lnSpc')
        spcPct = etree.SubElement(lnSpc, f'{{{ns_a}}}spcPct')
        spcPct.set('val', str(line_spacing_pct))

    if space_before is not None:
        spcBef = etree.SubElement(pPr, f'{{{ns_a}}}spcBef')
        spcPts = etree.SubElement(spcBef, f'{{{ns_a}}}spcPts')
        spcPts.set('val', str(space_before))

    if space_after is not None:
        spcAft = etree.SubElement(pPr, f'{{{ns_a}}}spcAft')
        spcPts = etree.SubElement(spcAft, f'{{{ns_a}}}spcPts')
        spcPts.set('val', str(space_after))


def _add_text_box(slide, left, top, width, height,
                  text, font_name=FONT_BODY, font_size=SIZE_BODY,
                  color=WHITE, bold=None, alignment=PP_ALIGN.LEFT,
                  anchor=MSO_ANCHOR.TOP, line_spacing=None):
    """단일 텍스트 박스 추가 (자간·행간 자동 적용)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    # 앵커 설정
    body_elem = tf._txBody
    body_pr = body_elem.find('{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
    if body_pr is not None:
        anchor_map = {
            MSO_ANCHOR.TOP: 't',
            MSO_ANCHOR.MIDDLE: 'ctr',
            MSO_ANCHOR.BOTTOM: 'b',
        }
        body_pr.set('anchor', anchor_map.get(anchor, 'ctr'))

    p = tf.paragraphs[0]
    p.alignment = alignment

    # 행간 적용
    ls = line_spacing if line_spacing else LINE_SPACING_NORMAL
    _apply_line_spacing(p, line_spacing_pct=ls,
                        space_before=SPACE_BEFORE, space_after=SPACE_AFTER)

    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.color.rgb = color
    if bold is not None:
        run.font.bold = bold

    # 자간 적용
    _apply_char_spacing(run, font_name)

    return txBox


def _parse_color(value, default=YELLOW):
    """색상 문자열/객체를 RGBColor로 정규화"""
    if isinstance(value, RGBColor):
        return value
    if not value:
        return default

    text = str(value).strip().upper().replace('#', '')
    named = {
        'YELLOW': YELLOW,
        'CYAN': CYAN,
        'RED': RED,
        'DARK_RED': RGBColor(0xC0, 0x00, 0x00),
        'WHITE': WHITE,
        'BLACK': BLACK,
    }
    if text in named:
        return named[text]
    if len(text) == 6:
        try:
            return RGBColor.from_string(text)
        except ValueError:
            return default
    return default


def _normalize_highlights(highlights, fallback_words=None):
    """기존 highlight[str]와 확장 highlights[dict]를 모두 수용"""
    normalized = []

    for item in highlights or []:
        if isinstance(item, str):
            normalized.append({
                'word': item,
                'color': YELLOW,
                'scale': 1.36,
                'bold': True,
            })
        elif isinstance(item, dict):
            word = item.get('word') or item.get('text')
            if not word:
                continue
            normalized.append({
                'word': word,
                'color': _parse_color(item.get('color'), YELLOW),
                'scale': float(item.get('scale', 1.36)),
                'bold': item.get('bold', True),
            })

    for item in fallback_words or []:
        if item not in {entry['word'] for entry in normalized}:
            normalized.append({
                'word': item,
                'color': YELLOW,
                'scale': 1.36,
                'bold': True,
            })

    return normalized


def _add_multiline_text(slide, left, top, width, height,
                        lines, font_name=FONT_BODY, font_size=SIZE_BODY,
                        color=WHITE, bold=None, alignment=PP_ALIGN.JUSTIFY,
                        highlights=None,
                        line_spacing=None):
    """여러 줄 텍스트 박스 (자간·행간 + 단어별 강조)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    ls = line_spacing if line_spacing else LINE_SPACING_NORMAL

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment

        # 행간·단락 간격 적용
        _apply_line_spacing(p, line_spacing_pct=ls,
                            space_before=SPACE_BEFORE, space_after=SPACE_AFTER)

        if highlights:
            _add_highlighted_runs(p, line, font_name, font_size, color,
                                  bold, highlights)
        else:
            run = p.add_run()
            run.text = line
            run.font.name = font_name
            run.font.size = font_size
            run.font.color.rgb = color
            if bold is not None:
                run.font.bold = bold
            # 자간 적용
            _apply_char_spacing(run, font_name)

    return txBox


def _add_highlighted_runs(paragraph, text, font_name, font_size, base_color,
                          bold, highlights):
    """텍스트 내 강조 단어를 색상+크기 증대로 적용"""
    remaining = text
    while remaining:
        earliest_pos = len(remaining)
        earliest = None
        for entry in highlights:
            word = entry['word']
            pos = remaining.find(word)
            if pos != -1 and pos < earliest_pos:
                earliest_pos = pos
                earliest = entry

        if earliest is None:
            run = paragraph.add_run()
            run.text = remaining
            run.font.name = font_name
            run.font.size = font_size
            run.font.color.rgb = base_color
            if bold is not None:
                run.font.bold = bold
            _apply_char_spacing(run, font_name)
            break
        else:
            if earliest_pos > 0:
                run = paragraph.add_run()
                run.text = remaining[:earliest_pos]
                run.font.name = font_name
                run.font.size = font_size
                run.font.color.rgb = base_color
                if bold is not None:
                    run.font.bold = bold
                _apply_char_spacing(run, font_name)

            run = paragraph.add_run()
            run.text = earliest['word']
            run.font.name = font_name
            run.font.size = Pt(font_size.pt * earliest.get('scale', 1.36))
            run.font.color.rgb = earliest.get('color', YELLOW)
            run.font.bold = earliest.get('bold', bold if bold is not None else True)
            _apply_char_spacing(run, font_name)

            remaining = remaining[earliest_pos + len(earliest['word']):]


def _add_header_bar(slide):
    """그라디언트 헤더 바 추가 (OXML 직접 삽입)"""
    # 고유 shape ID 생성
    shape_id = max(
        (int(sp.get('id', 0))
         for sp in slide._element.iter('{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr')),
        default=100
    ) + 1

    # 더 안전한 ID 생성: spTree 내 모든 cNvPr의 id 중 최대값 + 1
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    max_id = 0
    spTree = slide._element.find(f'{{{ns_p}}}cSld/{{{ns_p}}}spTree')
    if spTree is None:
        # fallback: p namespace 없이 시도
        spTree = slide._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}spTree')
    if spTree is not None:
        for cNvPr in spTree.iter('{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr'):
            try:
                max_id = max(max_id, int(cNvPr.get('id', 0)))
            except ValueError:
                pass

    shape_id = max(max_id + 1, 100)

    xml = HEADER_BAR_XML.format(shape_id=shape_id)
    sp_elem = etree.fromstring(xml.encode('utf-8'))

    # spTree에 삽입 (맨 뒤 = 맨 위 레이어)
    if spTree is not None:
        spTree.insert(2, sp_elem)  # nvGrpSpPr/grpSpPr(필수 2개) 다음 = 배경 레이어. insert(0)은 스키마 위반→PowerPoint 복구


def _add_image_fullscreen(slide, image_path):
    """풀스크린 이미지 추가 (맨 뒤 레이어)"""
    if not os.path.exists(image_path):
        print(f"  경고: 이미지 없음 → {image_path}")
        return None

    left, top, width, height = IMAGE_FULLSCREEN
    pic = slide.shapes.add_picture(image_path, left, top, width, height)

    # 맨 뒤로 보내기 (send to back)
    spTree = slide._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}spTree')
    if spTree is None:
        spTree = pic._element.getparent()
    sp_elem = pic._element
    spTree.remove(sp_elem)
    spTree.insert(2, sp_elem)  # nvGrpSpPr/grpSpPr 다음 = 배경 레이어. insert(0)은 스키마 위반→복구

    return pic


def _add_red_oval(slide, left, top, width, height):
    """빨간 타원 도형 추가 (마스크용, noFill + 9pt red line)"""
    oval = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, width, height
    )
    fill = oval.fill
    fill.background()

    oval.line.color.rgb = RED
    oval.line.width = Pt(9)

    return oval


def _add_red_rect(slide, left, top, width, height):
    """빨간 직사각형 추가 (마스크용, noFill + 9pt red line)"""
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    fill = rect.fill
    fill.background()

    rect.line.color.rgb = RED
    rect.line.width = Pt(9)

    return rect


def _add_scripture_frame(slide):
    left, top, width, height = SCRIPTURE_FRAME_POS
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.fill.background()
    rect.line.color.rgb = RED
    rect.line.width = Pt(9)
    return rect


def _add_scripture_header_tag(slide, ref_label):
    left, top, w, h = SCRIPTURE_TAG_POS
    tx_box = slide.shapes.add_textbox(left, top, w, h)
    tf = tx_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    _apply_line_spacing(p, line_spacing_pct=LINE_SPACING_NORMAL,
                        space_before=SPACE_BEFORE, space_after=SPACE_AFTER)

    run1 = p.add_run()
    run1.text = "성경 말씀"
    run1.font.name = FONT_HEADER
    run1.font.size = SIZE_HEADER_TAG
    run1.font.bold = True
    run1.font.color.rgb = WHITE
    _apply_char_spacing(run1, FONT_HEADER)

    run2 = p.add_run()
    run2.text = f"({ref_label})"
    run2.font.name = FONT_HEADER
    run2.font.size = Pt(32)
    run2.font.bold = True
    run2.font.color.rgb = WHITE
    _apply_char_spacing(run2, FONT_HEADER)
    return tx_box


# ============================================================
# 슬라이드 빌더 함수들
# ============================================================

def build_cover_scripture(prs, book_name, chapter_verse):
    """표지 1: "성경 말씀" + 책명(115pt 노랑) + 장절(80pt) + Fade In"""
    slide = _new_slide(prs)
    _add_header_bar(slide)

    # "성경 말씀" 태그
    left, top, w, h = COVER_TAG_POS
    tag_box = _add_text_box(slide, left, top, w, h,
                  "성경 말씀", FONT_HEADER, SIZE_HEADER_TAG, WHITE, bold=True)

    # 책 이름 (115pt 노랑) + 장절 (80pt)
    left, top, w, h = COVER_BOOK_POS
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]

    run_book = p.add_run()
    run_book.text = book_name + "\n"
    run_book.font.name = FONT_BODY
    run_book.font.size = SIZE_BOOK_NAME
    run_book.font.color.rgb = YELLOW
    _apply_char_spacing(run_book, FONT_BODY)

    run_cv = p.add_run()
    run_cv.text = chapter_verse
    run_cv.font.name = FONT_BODY
    run_cv.font.size = SIZE_CHAPTER_VERSE
    run_cv.font.color.rgb = WHITE
    _apply_char_spacing(run_cv, FONT_BODY)

    # 애니메이션: 태그 + 책명 동시 Fade In
    inject_multi_fade_in(slide, [tag_box.shape_id, txBox.shape_id])

    return slide


def build_cover_text(prs, ref_label, scripture_text):
    """표지 2: "성경 말씀(절)" + 본문 전문(60pt) + Fade In"""
    slide = _new_slide(prs)
    _add_header_bar(slide)

    tag_box = _add_scripture_header_tag(slide, ref_label)

    # 본문 텍스트
    left, top, w, h = SCRIPTURE_BOX_POS
    lines = scripture_text.split('\n') if '\n' in scripture_text else [scripture_text]
    text_box = _add_multiline_text(slide, left, top, w, h,
                        lines, FONT_BODY, SIZE_BODY, WHITE,
                        alignment=PP_ALIGN.JUSTIFY)

    # 애니메이션: 태그 + 본문 동시 Fade In
    inject_multi_fade_in(slide, [tag_box.shape_id, text_box.shape_id])

    return slide


def build_cover_title(prs, title, series_tag=None):
    """표지 3: "말씀 제목" + 설교 제목(96pt 노랑 Bold) + Fade In"""
    slide = _new_slide(prs)
    _add_header_bar(slide)

    shape_ids = []

    # "말씀 제목" 태그
    left, top, w, h = COVER_TAG_POS
    tag_box = _add_text_box(slide, left, top, w, h,
                  "말씀 제목", FONT_HEADER, SIZE_HEADER_TAG, WHITE, bold=True)
    shape_ids.append(tag_box.shape_id)

    # 시리즈 태그 (있으면)
    if series_tag:
        left, top, w, h = SERIES_TAG_POS
        st_box = _add_text_box(slide, left, top, w, h,
                      series_tag, FONT_BODY, Pt(45), YELLOW, bold=True)
        shape_ids.append(st_box.shape_id)

    # 설교 제목
    left, top, w, h = TITLE_POS
    title_box = _add_text_box(slide, left, top, w, h,
                  title, FONT_BODY, SIZE_TITLE, YELLOW,
                  bold=True, alignment=PP_ALIGN.LEFT)
    shape_ids.append(title_box.shape_id)

    # 애니메이션: 모든 요소 동시 Fade In
    inject_multi_fade_in(slide, shape_ids)

    return slide


def build_full_image(prs, image_path):
    """풀스크린 이미지 슬라이드 + Fade In"""
    slide = _new_slide(prs)
    pic = _add_image_fullscreen(slide, image_path)
    if pic:
        inject_fade_in(slide, pic.shape_id)
    return slide


def build_text_image(prs, image_path):
    """텍스트 PNG 풀스크린 + Fade In"""
    slide = _new_slide(prs)
    pic = _add_image_fullscreen(slide, image_path)
    if pic:
        inject_fade_in(slide, pic.shape_id)
    return slide


def build_question(prs, title_text=None, question_size=None, question_color=YELLOW,
                   image_path=None, question_position=None):
    """거대 ? 효과 슬라이드"""
    slide = _new_slide(prs)

    if question_size is None:
        question_size = SIZE_QUESTION_XL
    question_color = _parse_color(question_color, YELLOW)

    if image_path:
        _add_image_fullscreen(slide, image_path)

    # 제목 텍스트 (있으면)
    if title_text:
        left, top, w, h = TITLE_POS
        title_box = _add_text_box(slide, left, top, w, h,
                                  title_text, FONT_BODY, SIZE_TITLE, YELLOW,
                                  bold=True)

    # 빨간 타원
    oval = _add_red_oval(slide, Cm(16.23), Cm(4.85), Cm(10.49), Cm(9.99))
    oval_id = oval.shape_id

    size_pt = round(question_size.pt)
    q_height = QUESTION_HEIGHT_MAP.get(size_pt, QUESTION_POS[3])
    if question_position:
        left = Cm(question_position.get('left_cm', 18.72))
        top = Cm(question_position.get('top_cm', 8.83))
        w = Cm(question_position.get('width_cm', 8.54))
        h = Cm(question_position.get('height_cm', round(q_height / 360000, 2)))
    else:
        left, top, w, _ = QUESTION_POS
        h = q_height

    # 거대 ? 텍스트
    q_box = _add_text_box(slide, left, top, w, h,
                          "?", FONT_QUESTION, question_size, question_color,
                          alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text_id = q_box.shape_id

    # 애니메이션 주입
    inject_question_animation(slide, oval_id, text_id)

    return slide


def build_original_word(prs, word, pronunciation, meaning, font_size=None):
    """원어 제목 슬라이드 — 헬라어/히브리어 + 뜻 + 순차 Fade In"""
    slide = _new_slide(prs)

    if font_size is None:
        font_size = Pt(88)

    # 원어 (중앙 상단) — 클릭 1
    word_box = _add_text_box(slide, Cm(2), Cm(4), Cm(30), Cm(6),
                  word, FONT_GREEK, font_size, WHITE,
                  alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 음역 + 뜻 (중앙 하단) — 클릭 2
    sub_text = f"({pronunciation})\n{meaning}"
    meaning_box = _add_multiline_text(slide, Cm(2), Cm(11), Cm(30), Cm(6),
                        sub_text.split('\n'), FONT_BODY, SIZE_BODY, YELLOW,
                        alignment=PP_ALIGN.CENTER)

    # 애니메이션: 원어 먼저 Fade In → 클릭 → 뜻 Fade In
    inject_click_fade_in(slide, [word_box.shape_id, meaning_box.shape_id])

    return slide


def build_progressive_text(prs, lines, highlights=None, ref_label=None, masks=None):
    """점진 노출 — 본문 표시 + 특정 단어 강조 + Fade In"""
    slide = _new_slide(prs)

    shape_ids = []
    mask_ids = []
    highlights = _normalize_highlights(highlights)

    if ref_label:
        _add_header_bar(slide)
        frame = _add_scripture_frame(slide)
        tag_box = _add_scripture_header_tag(slide, ref_label)
        shape_ids.append(frame.shape_id)
        shape_ids.append(tag_box.shape_id)
        body_pos = SCRIPTURE_BOX_POS
    else:
        body_pos = (Cm(0.77), Cm(1.5), Cm(31.95), Cm(15.5))

    left, top, w, h = body_pos
    text_box = _add_multiline_text(slide, left, top, w, h,
                        lines, FONT_BODY, SIZE_BODY, WHITE,
                        alignment=PP_ALIGN.JUSTIFY,
                        highlights=highlights)
    shape_ids.append(text_box.shape_id)

    for mask in masks or []:
        shape_type = mask.get('type', 'rect')
        ml = mask.get('left_cm', 0)
        mt = mask.get('top_cm', 0)
        mw = mask.get('width_cm', 0)
        mh = mask.get('height_cm', 0)
        if shape_type == 'oval':
            shape = _add_red_oval(slide, Cm(ml), Cm(mt), Cm(mw), Cm(mh))
        else:
            shape = _add_red_rect(slide, Cm(ml), Cm(mt), Cm(mw), Cm(mh))
        mask_ids.append(shape.shape_id)

    if mask_ids:
        inject_text_mask_reveal(slide, mask_ids)
    else:
        inject_multi_fade_in(slide, shape_ids)

    return slide


def build_image_with_masks(prs, image_path, mask_positions):
    """이미지 + 빨간 마스크 리빌

    mask_positions: list of (type, left, top, width, height)
        type: 'oval' 또는 'rect'
        좌표: Cm 단위
    """
    slide = _new_slide(prs)
    _add_image_fullscreen(slide, image_path)

    mask_ids = []
    for mtype, ml, mt, mw, mh in mask_positions:
        if mtype == 'oval':
            shape = _add_red_oval(slide, Cm(ml), Cm(mt), Cm(mw), Cm(mh))
        else:
            shape = _add_red_rect(slide, Cm(ml), Cm(mt), Cm(mw), Cm(mh))
        mask_ids.append(shape.shape_id)

    # 멀티 마스크 애니메이션 주입
    if mask_ids:
        inject_multi_mask_animation(slide, mask_ids)

    return slide


def build_blank(prs):
    """빈 슬라이드"""
    return _new_slide(prs)
