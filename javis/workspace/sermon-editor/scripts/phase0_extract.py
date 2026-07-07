"""
Phase 0: 참조 PPT 전수 분석 및 XML 추출
- 60슬라이드 전수 분석 → slide_inventory.json
- 아키타입별 shape 구성 추출 → archetype_specs.json
- 애니메이션 XML 블록 추출 → anim_template_*.xml
- 그라디언트 바 XML 추출 → gradient_bar.xml
- 슬라이드 전환 XML 추출 → transition_fade.xml
- 템플릿 pptx 생성 → ppt-template-wide.pptx
"""

import sys
import os
import json
import copy

sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Emu, Pt, Inches
from lxml import etree

# === 경로 설정 ===
BASE = r"C:\Users\신이재\OneDrive - the presbyerian church of korea\바탕 화면\My-Sermon-Editor"
REF_PPT = os.path.join(BASE, "MYppt", "20260325(장년부수요기도회PPT)---wide.pptx")
OUTPUT_DIR = os.path.join(BASE, "scripts", "pptx_extracted")
TEMPLATE_DIR = os.path.join(BASE, "templates")

os.makedirs(OUTPUT_DIR, exist_ok=True)
os.makedirs(TEMPLATE_DIR, exist_ok=True)

# === 네임스페이스 ===
NSMAP = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'p14': 'http://schemas.microsoft.com/office/powerpoint/2010/main',
}


def emu_to_inches(emu):
    return round(emu / 914400, 2) if emu else 0


def extract_text_from_shape(shape):
    """shape에서 텍스트 추출 (줄바꿈 → |)"""
    if hasattr(shape, 'text') and shape.text.strip():
        return shape.text.strip().replace('\n', ' | ')[:150]
    return ""


def extract_font_info(shape):
    """shape에서 폰트 정보 추출"""
    fonts = set()
    sizes = set()
    colors = set()
    try:
        if hasattr(shape, 'text_frame'):
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.name:
                        fonts.add(run.font.name)
                    if run.font.size:
                        sizes.add(int(run.font.size / 12700))  # EMU to pt
                    if run.font.color and run.font.color.rgb:
                        colors.add(str(run.font.color.rgb))
    except:
        pass
    return list(fonts), list(sizes), list(colors)


def analyze_slide(slide, idx):
    """단일 슬라이드 전수 분석"""
    info = {
        "slide_num": idx + 1,
        "layout": slide.slide_layout.name,
        "shapes": [],
        "has_image": False,
        "has_animation": False,
        "has_transition": False,
        "texts": [],
        "images": [],
        "animation_xml_size": 0,
    }

    # 애니메이션 확인
    ns_p = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
    timing = slide.element.find(f'.//{ns_p}timing')
    if timing is not None:
        info["has_animation"] = True
        info["animation_xml_size"] = len(etree.tostring(timing))

    # 전환 확인
    transition = slide.element.find(f'{ns_p}transition')
    if transition is not None:
        info["has_transition"] = True

    # shape 분석
    for shape in slide.shapes:
        shape_info = {
            "name": shape.name,
            "type": str(shape.shape_type),
            "left_in": emu_to_inches(shape.left),
            "top_in": emu_to_inches(shape.top),
            "width_in": emu_to_inches(shape.width),
            "height_in": emu_to_inches(shape.height),
            "shape_id": shape.shape_id,
        }

        # 이미지
        if shape.shape_type == 13:  # PICTURE
            info["has_image"] = True
            info["images"].append({
                "width_in": emu_to_inches(shape.width),
                "height_in": emu_to_inches(shape.height),
                "left_in": emu_to_inches(shape.left),
                "top_in": emu_to_inches(shape.top),
            })

        # 텍스트
        text = extract_text_from_shape(shape)
        if text:
            fonts, sizes, colors = extract_font_info(shape)
            shape_info["text"] = text
            shape_info["fonts"] = fonts
            shape_info["sizes"] = sizes
            shape_info["colors"] = colors
            info["texts"].append(text)

        # 그라디언트 감지
        fill_elem = shape.element.find('.//' + '{http://schemas.openxmlformats.org/drawingml/2006/main}gradFill')
        if fill_elem is not None:
            shape_info["has_gradient"] = True

        info["shapes"].append(shape_info)

    return info


def classify_slide_type(info):
    """슬라이드 타입 자동 분류"""
    texts = " ".join(info["texts"]).lower()
    has_img = info["has_image"]
    shape_count = len(info["shapes"])
    layout = info["layout"]

    # 빈 슬라이드
    if not info["texts"] and not has_img:
        return "blank"

    # 찬양/악보 (마지막 구간, 빈 화면 레이아웃 + 이미지만)
    if layout == '빈 화면' and has_img and not info["texts"]:
        return "hymn"

    # 기도 슬라이드
    if "기도" in texts and shape_count <= 3:
        return "prayer"

    # 물음표
    if "?" in texts and len(texts) < 5:
        return "question"

    # 성경 구절 (그라디언트 바 존재 + 성경 텍스트)
    has_gradient = any(s.get("has_gradient") for s in info["shapes"])
    if has_gradient and info["texts"]:
        return "scripture"

    # 이미지 + 텍스트
    if has_img and info["texts"]:
        # 텍스트가 짧으면 버블, 길면 오버레이
        total_text_len = sum(len(t) for t in info["texts"])
        if total_text_len < 30:
            return "image_bubble"
        else:
            return "image_text"

    # 이미지만
    if has_img and not info["texts"]:
        return "full_image"

    # 텍스트만
    if info["texts"] and not has_img:
        return "text_only"

    return "unknown"


def extract_animation_xml(slide, slide_num, output_dir):
    """슬라이드에서 애니메이션 XML 블록 추출"""
    ns_p = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
    timing = slide.element.find(f'.//{ns_p}timing')
    if timing is not None:
        xml_str = etree.tostring(timing, pretty_print=True, encoding='unicode')
        filepath = os.path.join(output_dir, f"anim_slide_{slide_num:02d}.xml")
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(xml_str)
        return filepath
    return None


def extract_transition_xml(slide, output_dir):
    """슬라이드 전환 XML 추출"""
    ns_p = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
    transition = slide.element.find(f'{ns_p}transition')
    if transition is not None:
        xml_str = etree.tostring(transition, pretty_print=True, encoding='unicode')
        filepath = os.path.join(output_dir, "transition_fade.xml")
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(xml_str)
        return True
    return False


def extract_gradient_xml(slide, output_dir):
    """그라디언트 바 shape의 XML 추출"""
    ns_a = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
    for shape in slide.shapes:
        grad = shape.element.find(f'.//{ns_a}gradFill')
        if grad is not None:
            # shape 전체 XML 추출
            shape_xml = etree.tostring(shape.element, pretty_print=True, encoding='unicode')
            filepath = os.path.join(output_dir, "gradient_bar_shape.xml")
            with open(filepath, 'w', encoding='utf-8') as f:
                f.write(shape_xml)
            # gradFill만 추출
            grad_xml = etree.tostring(grad, pretty_print=True, encoding='unicode')
            filepath2 = os.path.join(output_dir, "gradient_fill.xml")
            with open(filepath2, 'w', encoding='utf-8') as f:
                f.write(grad_xml)
            return True
    return False


def create_template(ref_ppt_path, template_path):
    """참조 PPT에서 슬라이드를 제거하고 마스터/레이아웃만 보존한 템플릿 생성"""
    prs = Presentation(ref_ppt_path)

    # 모든 슬라이드 ID 수집
    sldIdLst = prs.slides._sldIdLst
    rIds = [sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
            for sldId in sldIdLst]

    # 역순으로 삭제 (인덱스 안정성)
    for rId in reversed(rIds):
        try:
            prs.part.drop_rel(rId)
        except:
            pass
    # sldIdLst 비우기
    for child in list(sldIdLst):
        sldIdLst.remove(child)

    prs.save(template_path)
    print(f"템플릿 생성: {template_path}")


def main():
    print(f"참조 PPT: {REF_PPT}")
    print(f"출력 디렉토리: {OUTPUT_DIR}")

    prs = Presentation(REF_PPT)
    print(f"슬라이드 크기: {emu_to_inches(prs.slide_width)}x{emu_to_inches(prs.slide_height)} inches")
    print(f"총 슬라이드: {len(prs.slides)}장")

    # === 0-1: 60장 전수 분석 ===
    print("\n=== Phase 0-1: 전수 분석 ===")
    inventory = []
    transition_extracted = False
    gradient_extracted = False

    for i, slide in enumerate(prs.slides):
        info = analyze_slide(slide, i)
        info["slide_type"] = classify_slide_type(info)

        # 애니메이션 XML 추출 (0-5)
        if info["has_animation"]:
            anim_path = extract_animation_xml(slide, i + 1, OUTPUT_DIR)
            if anim_path:
                info["animation_xml_file"] = os.path.basename(anim_path)

        # 전환 XML 추출 (0-5, 첫 번째만)
        if info["has_transition"] and not transition_extracted:
            if extract_transition_xml(slide, OUTPUT_DIR):
                transition_extracted = True

        # 그라디언트 XML 추출 (0-6, 첫 번째만)
        if not gradient_extracted:
            if extract_gradient_xml(slide, OUTPUT_DIR):
                gradient_extracted = True

        inventory.append(info)

        # 콘솔 출력
        type_str = info["slide_type"].upper()
        anim_str = "[ANIM]" if info["has_animation"] else ""
        img_str = "[IMG]" if info["has_image"] else ""
        text_preview = info["texts"][0][:60] if info["texts"] else ""
        print(f"  S{i+1:02d} {type_str:15s} {img_str:6s} {anim_str:7s} {text_preview}")

    # 인벤토리 저장
    inv_path = os.path.join(OUTPUT_DIR, "slide_inventory.json")
    with open(inv_path, 'w', encoding='utf-8') as f:
        json.dump(inventory, f, ensure_ascii=False, indent=2)
    print(f"\n인벤토리 저장: {inv_path}")

    # === 0-2: 타입별 통계 ===
    print("\n=== Phase 0-2: 타입별 통계 ===")
    type_counts = {}
    for info in inventory:
        t = info["slide_type"]
        type_counts[t] = type_counts.get(t, 0) + 1
    for t, c in sorted(type_counts.items(), key=lambda x: -x[1]):
        print(f"  {t:15s}: {c}장")

    # === 0-3: 아키타입 슬라이드 선별 ===
    print("\n=== Phase 0-3: 아키타입 선별 ===")
    archetypes = {}
    for info in inventory:
        t = info["slide_type"]
        if t not in archetypes:
            archetypes[t] = info["slide_num"]
    arch_path = os.path.join(OUTPUT_DIR, "archetypes.json")
    with open(arch_path, 'w', encoding='utf-8') as f:
        json.dump(archetypes, f, ensure_ascii=False, indent=2)
    for t, num in archetypes.items():
        print(f"  {t:15s} → S{num:02d}")

    # === 0-4: 아키타입별 shape 상세 추출 ===
    print("\n=== Phase 0-4: 아키타입 shape 상세 ===")
    archetype_details = {}
    for t, slide_num in archetypes.items():
        slide_info = inventory[slide_num - 1]
        archetype_details[t] = {
            "slide_num": slide_num,
            "layout": slide_info["layout"],
            "shape_count": len(slide_info["shapes"]),
            "shapes": slide_info["shapes"],
        }
    details_path = os.path.join(OUTPUT_DIR, "archetype_specs.json")
    with open(details_path, 'w', encoding='utf-8') as f:
        json.dump(archetype_details, f, ensure_ascii=False, indent=2)
    print(f"아키타입 상세 저장: {details_path}")

    # === 0-7: 템플릿 생성 ===
    print("\n=== Phase 0-7: 템플릿 생성 ===")
    template_path = os.path.join(TEMPLATE_DIR, "ppt-template-wide.pptx")
    try:
        create_template(REF_PPT, template_path)
    except Exception as e:
        print(f"템플릿 생성 실패: {e}")
        print("대안: 참조 PPT를 직접 복사하여 수동으로 슬라이드 삭제 필요")

    # === 요약 ===
    print("\n=== Phase 0 완료 요약 ===")
    extracted_files = os.listdir(OUTPUT_DIR)
    print(f"추출 파일 수: {len(extracted_files)}")
    for f in sorted(extracted_files):
        size = os.path.getsize(os.path.join(OUTPUT_DIR, f))
        print(f"  {f} ({size:,} bytes)")

    anim_files = [f for f in extracted_files if f.startswith("anim_slide_")]
    print(f"\n애니메이션 XML 파일: {len(anim_files)}개")
    print(f"전환 XML: {'있음' if transition_extracted else '없음'}")
    print(f"그라디언트 XML: {'있음' if gradient_extracted else '없음'}")
    print(f"템플릿: {'생성됨' if os.path.exists(template_path) else '실패'}")


if __name__ == "__main__":
    main()
