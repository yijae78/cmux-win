"""본문 슬라이드 = 수정 가능한 네이티브 텍스트 (미리보기 이미지와 동일 디자인)

설계 (2026-06-06 신교수님 지시):
- 미리보기 HTML = 지금처럼 text_PP.png 이미지 그대로 (md_preview.py 무변경)
- PPT = 본문 슬롯을 '수정 가능한 네이티브 텍스트'로. 단 디자인·형식·글씨체는
  generate_text_slides.py(미리보기 이미지)와 똑같이.

재현 전략:
- 캔버스 1280x720 ↔ PPT 13.33"x7.5" 환산이 정확히 떨어진다: 1px = 9525 EMU = 0.75pt
- generate_text_slides.py의 레이아웃 계산(헤더 밴드 높이·제목/본문 폰트 자동축소·
  줄바꿈·줄간격·세로중앙)을 그대로 재사용 → 위치·크기 픽셀 일치
- 한글=나눔명조 ExtraBold(ea 폰트 지정), 헬라/히브리=Times New Roman Bold(latin),
  한자=함초롬바탕. 강조 구절=노랑 run.

공용: sermon_pptx.py (전 부서 공용 메인 엔진)에서 호출.

단독 테스트:
    python scripts/pptx_text_native.py <설교폴더> <슬롯번호>   # 1장짜리 PPT 생성
"""

import os
import sys
from pathlib import Path

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from PIL import Image, ImageDraw

# 미리보기 렌더러의 레이아웃 계산·폰트 헬퍼를 그대로 재사용 (디자인 단일 출처)
import generate_text_slides as G

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN, MSO_AUTO_SIZE
from pptx.oxml.ns import qn


# ── 환산 상수 (정확히 떨어짐) ───────────────────────────────
EMU_PER_PX = 9525            # 12192000 EMU / 1280 px
PT_PER_PX = 0.75             # 960 pt / 1280 px

# ── 폰트 이름 (PowerPoint에 설치된 이름) ────────────────────
KOR_FONT = "나눔명조 ExtraBold"
ORIG_FONT = "Times New Roman"      # 헬라어/히브리어 (bold)
CJK_FONT = "함초롬바탕"

# ── 색 (generate_text_slides.py와 동일) ─────────────────────
C_BG = RGBColor(10, 14, 26)
C_BAND = RGBColor(14, 32, 62)
C_GOLD = RGBColor(251, 191, 36)
C_TITLE = RGBColor(255, 255, 255)
C_BODY = RGBColor(241, 245, 249)
C_HL = RGBColor(255, 255, 0)

# 줄별 박스를 미리보기 글자 top(cy)에 픽셀 일치시키는 세로 보정 계수.
# box_top = cy - body_size*K. PowerPoint(잉크 offset)와 PIL(draw offset)의 차이를
# 상쇄하도록 측정으로 확정: 음수 = box를 cy보다 살짝 아래로. (PP05/15/17 3슬롯 회귀)
BODY_LEAD_K = -0.107


def px_emu(px):
    return Emu(int(round(px * EMU_PER_PX)))


def px_pt(px):
    return Pt(px * PT_PER_PX)


def _set_run_font(run, font_name, size_px, color, bold=True):
    """run에 폰트(latin+ea+cs 모두) / 크기 / 색 / bold 지정.

    한글은 ea(East Asian) 폰트로 지정해야 PowerPoint가 나눔명조를 적용한다.
    latin만 지정하면 한글이 테마 기본 글꼴로 깨진다.
    """
    f = run.font
    f.size = px_pt(size_px)
    f.bold = bold
    f.color.rgb = color
    f.name = font_name  # <a:latin> 설정
    rPr = run._r.get_or_add_rPr()
    # latin 다음 위치에 ea, cs 삽입 (스키마 순서 보존)
    latin = rPr.find(qn('a:latin'))
    for tag in ('a:cs', 'a:ea'):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            if latin is not None:
                latin.addnext(el)
            else:
                rPr.append(el)
        el.set('typeface', font_name)


def _font_name_for_cat(cat):
    if cat == 'orig':
        return ORIG_FONT
    if cat == 'cjk':
        return CJK_FONT
    return KOR_FONT


def _line_runs(line, highlight, hl_color, base_color):
    """generate_text_slides.draw_mixed_line과 동일한 분해:
    1) 강조 단어로 색 청크 분할 → 2) 각 청크를 원어/한자/한글 run으로 분할.
    반환: [(run_text, font_cat, color), ...]
    """
    sorted_hl = sorted(highlight or [], key=len, reverse=True)
    chunks = []
    rem = line
    while rem:
        hit = None
        for hw in sorted_hl:
            if hw and rem.startswith(hw):
                hit = hw
                break
        if hit:
            chunks.append((hit, hl_color))
            rem = rem[len(hit):]
        else:
            chunks.append((rem[0], base_color))
            rem = rem[1:]
    # 같은 색 합치기
    merged = []
    for tok, col in chunks:
        if merged and merged[-1][1] == col:
            merged[-1][0] += tok
        else:
            merged.append([tok, col])
    # 폰트 카테고리로 재분할
    out = []
    for tok, col in merged:
        for run_text, cat in G.split_runs(tok):
            out.append((run_text, cat, col))
    return out


def compute_layout(title, body):
    """generate_text_slides.generate_text_slide의 레이아웃 계산만 추출.

    반환 dict: band_h, title_size, title_lines, title_line_h,
               body_size, wrapped_lines, line_h, body_top, start_y, body_max_h
    """
    dummy = Image.new("RGB", (G.W, G.H))
    draw = ImageDraw.Draw(dummy)

    title_lines = [l.strip() for l in title.split("\n") if l.strip()]

    # 제목 폰트 자동축소 (헤더 밴드 기준: 80 → 44)
    title_size = 80
    while title_size >= 44:
        tk = G.find_font(G.KOREAN_FONT_CANDIDATES, title_size)
        to = G.find_font(G.ORIGINAL_LANG_FONT_CANDIDATES, title_size)
        tc = G.find_font(G.CJK_FONT_CANDIDATES, title_size)
        max_tw = max(
            (G.measure_line_width(draw, l, tk, to, tc) for l in title_lines),
            default=0,
        )
        if max_tw <= G.W - G.LEFT_MARGIN - 40:
            break
        title_size -= 4
    tk = G.find_font(G.KOREAN_FONT_CANDIDATES, title_size)
    to = G.find_font(G.ORIGINAL_LANG_FONT_CANDIDATES, title_size)
    title_line_h = max(G.text_height(tk), G.text_height(to)) + 6
    title_y = 35
    band_h = title_y + len(title_lines) * title_line_h + 22

    # 본문 영역
    body_top = band_h + 30
    body_max_w = G.W - G.LEFT_MARGIN - 60
    body_max_h = G.H - body_top - 50

    body_lines = body.split("\n")
    nonblank = sum(1 for l in body_lines if l.strip())
    line_pad = 18 if nonblank <= 4 else (14 if nonblank <= 6 else 10)

    body_size = 80
    wrapped_lines = body_lines
    line_h = 0
    while body_size >= 44:
        bk = G.find_font(G.KOREAN_FONT_CANDIDATES, body_size)
        bo = G.find_font(G.ORIGINAL_LANG_FONT_CANDIDATES, body_size)
        bc = G.find_font(G.CJK_FONT_CANDIDATES, body_size)
        line_h = max(G.text_height(bk), G.text_height(bo)) + line_pad
        wrapped_lines = []
        for line in body_lines:
            if not line.strip():
                wrapped_lines.append("")
            else:
                wrapped_lines.extend(
                    G.wrap_line(draw, line.strip(), bk, bo, bc, body_max_w)
                )
        total_h = sum(line_h if l.strip() else line_h // 2 for l in wrapped_lines)
        if total_h <= body_max_h:
            break
        body_size -= 4

    total_h = sum(line_h if l.strip() else line_h // 2 for l in wrapped_lines)
    start_y = body_top + max(0, (body_max_h - total_h) // 2)

    return {
        "band_h": band_h,
        "title_size": title_size,
        "title_lines": title_lines,
        "title_line_h": title_line_h,
        "title_y": title_y,
        "body_size": body_size,
        "wrapped_lines": wrapped_lines,
        "line_h": line_h,
        "body_top": body_top,
        "start_y": start_y,
        "total_h": total_h,
        "body_max_h": body_max_h,
    }


def _add_rect(slide, x_px, y_px, w_px, h_px, color):
    from pptx.enum.shapes import MSO_SHAPE
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                px_emu(x_px), px_emu(y_px), px_emu(w_px), px_emu(h_px))
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def build_native_text_slide(prs, title, body, highlight):
    """본문 슬롯 1장을 네이티브 텍스트로 추가. 미리보기 이미지와 동일 디자인."""
    if isinstance(highlight, str):
        highlight = [highlight] if highlight.strip() else []

    L = compute_layout(title, body)

    # blank 레이아웃 (없으면 마지막)
    layout = None
    for lay in prs.slide_layouts:
        if (lay.name or "").lower() in ("blank", "빈 화면", "빈화면"):
            layout = lay
            break
    if layout is None:
        layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(layout)

    # 1) 배경 (전체 검정)
    _add_rect(slide, 0, 0, G.W, G.H, C_BG)
    # 2) 헤더 밴드 (진네이비) + 황금 하단 엣지 6px
    _add_rect(slide, 0, 0, G.W, L["band_h"], C_BAND)
    _add_rect(slide, 0, L["band_h"] - 6, G.W, 6, C_GOLD)
    # 3) 좌측 황금 세로 바 (0~10px)
    _add_rect(slide, 0, 0, 10, G.H, C_GOLD)

    # 4) 제목 — 글자 잉크의 세로 중심을 헤더 밴드의 정중앙에 맞춘다.
    #    PowerPoint 텍스트박스는 box_top에서 글자 잉크 중심까지 title_size*INK_CENTER_K 아래.
    #    (렌더 측정 회귀 — PP05: box_top=37.5 → 잉크중심=85 → offset/size = 47.5/80 = 0.594)
    #    title_lead(본문용 보정)는 제목에 쓰지 않는다 — 그게 글자를 아래로 밀던 원인.
    #    신교수님 지시 2026-06-06: 밴드 위아래 정확히 가운데로.
    INK_CENTER_K = 0.600
    n_title = max(1, len(L["title_lines"]))
    ty = L["band_h"] / 2.0 - L["title_size"] * INK_CENTER_K - (n_title - 1) * L["title_line_h"] / 2.0
    for tl in L["title_lines"]:
        _add_text_line(slide, G.LEFT_MARGIN, ty,
                       G.W - G.LEFT_MARGIN - 40, L["title_line_h"],
                       tl, [], L["title_size"], C_TITLE)
        ty += L["title_line_h"]

    # 5) 본문 — 줄별 고정 박스. 각 줄을 미리보기 cy 좌표에 못 박아 픽셀 일치시킨다.
    #    (한 박스에 본문 전체를 넣으면 PowerPoint의 line_spacing이 폰트 metrics에 따라
    #     비선형으로 흔들려 슬롯마다 ±6px 어긋남 → 줄별 박스로 그 변수를 제거.)
    #    BODY_LEAD: PowerPoint 텍스트박스는 글자를 box_top보다 약간 아래에 그리므로
    #    그만큼 박스를 위로 올려 미리보기 글자 top(cy)에 일치시킨다 (폰트 비례).
    BODY_LEAD = round(L["body_size"] * BODY_LEAD_K)
    cy = L["start_y"]
    for line in L["wrapped_lines"]:
        if not line.strip():
            cy += L["line_h"] // 2
            continue
        _add_text_line(slide, G.LEFT_MARGIN, cy - BODY_LEAD,
                       G.W - G.LEFT_MARGIN - 60, L["line_h"],
                       line.strip(), highlight, L["body_size"], C_BODY)
        cy += L["line_h"]

    return slide


def _add_text_line(slide, x_px, y_px, w_px, h_px, line, highlight, size_px, base_color):
    """한 줄 = 한 텍스트박스 (위치 고정). run별 폰트/색 전환."""
    tb = slide.shapes.add_textbox(px_emu(x_px), px_emu(y_px), px_emu(w_px), px_emu(h_px))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.line_spacing = 1.0
    for run_text, cat, col in _line_runs(line, highlight, C_HL, base_color):
        r = p.add_run()
        r.text = run_text
        _set_run_font(r, _font_name_for_cat(cat), size_px, col)
    return tb


# ── 단독 테스트: 1장짜리 PPT 생성 ──────────────────────────
def _make_single(sermon_dir, slot_num):
    import json
    sermon_dir = Path(sermon_dir).resolve()
    slots = json.loads((sermon_dir / "slots.json").read_text(encoding="utf-8"))
    target = None
    for s in slots:
        if str(s.get("num")) == str(slot_num):
            target = s
            break
    if target is None:
        print(f"슬롯 {slot_num} 없음")
        sys.exit(1)
    t = target.get("text", {}) or {}
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    build_native_text_slide(prs, t.get("title", ""), t.get("body", ""), t.get("highlight", []))
    out = sermon_dir / "output" / f"_sample_native_PP{slot_num}.pptx"
    out.parent.mkdir(exist_ok=True)
    prs.save(str(out))
    print(f"[OK] {out}")
    return str(out)


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("사용: python scripts/pptx_text_native.py <설교폴더> <슬롯번호>")
        sys.exit(1)
    _make_single(sys.argv[1], sys.argv[2])
