"""
설교 PPT 자동 생성 — 전 부서 공용 메인 엔진 (장년부·아동부·새벽 등)

⭐ 지금 형식 표준 (2026-06-06 「우산」 정립):
- 본문 슬롯 = 수정 가능한 네이티브 텍스트 (pptx_text_native)
  → 미리보기 이미지와 동일 디자인: 깊은 검정 배경 + 좌측 황금 바
    + 진네이비 헤더 밴드(글씨 세로 정중앙) + 황금 하단 엣지 + 노랑 강조
- 이미지 슬롯 = NLM 이미지 풀스크린
- 표지 3장(성경말씀/본문/제목)도 동일 밴드 디자인으로 통일
- 글꼴 임베딩은 스킬(ppt-maker)에서 PowerPoint COM으로 별도 적용

사용법:
    python scripts/sermon_pptx.py <설교폴더>

입력 (부서 자동 감지 — kids/ 우선, 없으면 직접):
    - slots.json (또는 kids/slots.json) — 슬롯 메타데이터
    - images/ppt_NN.png (또는 kids_ppt_NN.png) — 이미지 슬롯
    - 본문 슬롯은 slots.json의 text 필드에서 네이티브 렌더 (PNG 불필요)

출력:
    - output/{부서-}{제목}.pptx
"""

import sys
import os
import json
import re
from pathlib import Path

sys.stdout.reconfigure(encoding='utf-8')
sys.stderr.reconfigure(encoding='utf-8')

# scripts/ 디렉토리를 path에 추가
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from pptx import Presentation
from pptx.util import Pt

from kids_pptx_design import SLIDE_WIDTH, SLIDE_HEIGHT
from kids_pptx_builder import (
    build_cover_scripture,
    build_cover_text,
    build_cover_title,
    build_full_image,
    build_text_image,
    build_question,
    build_original_word,
    build_progressive_text,
    build_image_with_masks,
    build_blank,
)
# 본문 슬롯 = 수정 가능한 네이티브 텍스트 (미리보기 이미지와 동일 디자인, 픽셀 일치)
from pptx_text_native import build_native_text_slide


def _is_original_word_slot(slot):
    """원어(헬라어/히브리어) 슬롯인지 감지
    판별 기준: title에 헬라어/히브리어 문자가 있거나, desc에 '원어' 키워드
    주의: desc에 헬라어가 있더라도 '원어' 키워드 없으면 일반 텍스트 슬롯
    """
    text_info = slot.get('text', {})
    title = text_info.get('title', '') if text_info else ''
    desc = slot.get('desc', '')
    # 헬라어(기본+확장) 또는 히브리어 문자가 제목에 포함
    if re.search(r'[\u0370-\u03FF\u1F00-\u1FFF\u0590-\u05FF]', title):
        return True
    # desc에 "원어" 키워드 명시
    if '원어' in desc:
        return True
    return False


def _extract_original_word(slot):
    """원어 슬롯에서 word, pronunciation, meaning 추출"""
    text_info = slot.get('text', {})
    title = text_info.get('title', '') if text_info else ''
    body = text_info.get('body', '') if text_info else ''
    desc = slot.get('desc', '')

    # title에서 헬라어/히브리어 추출 (확장 범위 포함)
    m = re.search(r'[\u0370-\u03FF\u1F00-\u1FFF\u0590-\u05FF]+', title)
    if m:
        word = m.group(0)
    else:
        # desc에서 추출 시도
        m = re.search(r'[\u0370-\u03FF\u1F00-\u1FFF\u0590-\u05FF]+', desc)
        word = m.group(0) if m else title

    # body에서 뜻 추출
    meaning = body.replace('\n', '\n')  # 줄바꿈 유지
    pronunciation = ''

    return word, pronunciation, meaning


def _looks_like_scripture_ref(text):
    """텍스트 슬롯 title이 성경 참조인지 대략 판별"""
    if not text:
        return False
    return bool(
        re.match(r'^.+?\s*\d+:\d+', text)
        or re.match(r'^[가-힣A-Za-z]+\s*\d+:\d+', text)
    )


def _extract_highlights(text_info):
    """기존 highlight와 확장 highlights를 모두 수용"""
    if not text_info:
        return []
    if text_info.get('highlights'):
        return text_info.get('highlights', [])
    return text_info.get('highlight', [])


def _extract_masks(slot):
    """slot/text 하위 masks 메타 수용"""
    text_info = slot.get('text', {}) if isinstance(slot.get('text'), dict) else {}
    return slot.get('masks') or text_info.get('masks') or []


def _is_question_slot(slot):
    """거대 물음표 전용 슬롯 판별"""
    slot_type = str(slot.get('type', '')).lower()
    if slot_type == 'question':
        return True
    text_info = slot.get('text', {}) if isinstance(slot.get('text'), dict) else {}
    if str(text_info.get('type', '')).lower() == 'question':
        return True
    title = text_info.get('title', '')
    body = text_info.get('body', '')
    return title.strip() == '?' or body.strip() == '?'


def load_slots(sermon_dir):
    """kids/slots.json 로드"""
    slots_path = sermon_dir / "kids" / "slots.json"
    if not slots_path.exists():
        # fallback: sermon_dir 직접
        slots_path = sermon_dir / "slots.json"
    if not slots_path.exists():
        print(f"오류: slots.json을 찾을 수 없습니다: {slots_path}")
        sys.exit(1)

    with open(slots_path, 'r', encoding='utf-8') as f:
        return json.load(f)


def extract_metadata(slots):
    """slots.json에서 표지 정보 추출 (첫 슬롯 = 표지)"""
    if not slots:
        return {}, []

    first = slots[0]
    meta = {}

    # 표지 슬롯에서 성경 참조, 제목 추출
    if first.get('text'):
        _t = first['text'].get('title', '')
        _b = first['text'].get('body', '')
        # 미리보기용 PP01 표지 라벨("본문 :", "제목 :") 제거
        meta['book_ref'] = re.sub(r'^\s*본문\s*[:：]\s*', '', _t).strip()
        meta['title'] = re.sub(r'^\s*제목\s*[:：]\s*', '', _b).strip()

    # desc에서 추출 시도
    if not meta.get('book_ref') and first.get('desc'):
        desc = first['desc']
        # "표지 — 고린도전서 4:1-5 「칭찬」" 패턴
        m = re.search(r'[—-]\s*(.+?)\s*[「「](.+?)[」」]', desc)
        if m:
            meta['book_ref'] = m.group(1).strip()
            meta['title'] = m.group(2).strip()
        else:
            # "표지 — 고린도전서 4:1-5" 패턴
            m = re.search(r'[—-]\s*(.+)', desc)
            if m:
                meta['book_ref'] = m.group(1).strip()

    # 성경 책 이름과 장절 분리
    book_ref = meta.get('book_ref', '')
    # "고린도전서 4:1-5" → 책="고린도전서", 절="4장 1-5절"
    m = re.match(r'^(.+?)\s*(\d+):(.+)$', book_ref)
    if m:
        meta['book_name'] = m.group(1)
        meta['chapter_verse'] = f"{m.group(2)}장 {m.group(3)}절"
        meta['ref_label'] = f"{meta['book_name'][:1]}{m.group(2)}:{m.group(3)}"
    else:
        m = re.match(r'^(.+?)\s*(\d+)장\s*(.+)$', book_ref)
        if m:
            meta['book_name'] = m.group(1)
            meta['chapter_verse'] = f"{m.group(2)}장 {m.group(3)}"
            meta['ref_label'] = book_ref
        else:
            meta['book_name'] = book_ref
            meta['chapter_verse'] = ''
            meta['ref_label'] = book_ref

    # 본문 슬롯 (첫 슬롯 제외)
    content_slots = slots[1:]

    return meta, content_slots


def find_image(sermon_dir, slot_num, is_image_slot):
    """슬롯에 대응하는 이미지 파일 찾기"""
    images_dir = sermon_dir / "images"

    if is_image_slot:
        # NLM 이미지: kids_ppt_NN.png
        candidates = [
            images_dir / f"kids_ppt_{slot_num}.png",
            images_dir / f"kids_ppt_{slot_num}.jpg",
            images_dir / f"ppt_{slot_num}.png",
        ]
    else:
        # 텍스트 슬라이드: kids_text_PPNN.png
        candidates = [
            images_dir / f"kids_text_PP{slot_num}.png",
            images_dir / f"text_PP{slot_num}.png",
        ]
        # kids/ 하위 폴더도 탐색
        kids_images = sermon_dir / "kids" / "images"
        if kids_images.exists():
            candidates.extend([
                kids_images / f"text_PP{slot_num}.png",
            ])

    for c in candidates:
        if c.exists():
            return str(c)

    return None


def find_scripture_text(sermon_dir):
    """원고에서 성경 본문 텍스트 추출 (표지 2용)"""
    for name in ["원고-아동부.md", "원고.md"]:
        manuscript = sermon_dir / name
        if manuscript.exists():
            with open(manuscript, 'r', encoding='utf-8') as f:
                content = f.read()
            # ```로 감싼 첫 번째 코드 블록을 성경 본문으로 추정
            m = re.search(r'```\n(.*?)\n```', content, re.DOTALL)
            if m:
                return m.group(1).strip()
    return None


def generate_kids_pptx(sermon_dir, template_path=None):
    """아동부 PPT 생성 메인 함수"""
    sermon_dir = Path(sermon_dir).resolve()
    output_dir = sermon_dir / "output"
    output_dir.mkdir(exist_ok=True)

    print("=" * 60)
    print("아동부 PPT 자동 생성")
    print("=" * 60)
    print(f"설교 폴더: {sermon_dir}")

    # 1. slots.json 로드
    print("\n[1/4] slots.json 로드 중...")
    slots = load_slots(sermon_dir)
    meta, content_slots = extract_metadata(slots)

    book_name = meta.get('book_name', '성경')
    chapter_verse = meta.get('chapter_verse', '')
    title = meta.get('title', '제목')
    ref_label = meta.get('ref_label', '')

    print(f"  성경: {book_name} {chapter_verse}")
    print(f"  제목: {title}")
    print(f"  슬롯: {len(slots)}개 (표지 1 + 본문 {len(content_slots)})")

    # 2. 프레젠테이션 생성
    print("\n[2/4] 프레젠테이션 생성 중...")

    # 템플릿 탐색
    if template_path is None:
        project_root = sermon_dir.parent
        if project_root.name == 'sermons':
            project_root = project_root.parent
        tp = project_root / "templates" / "ppt-template-wide.pptx"
        if tp.exists():
            template_path = str(tp)

    if template_path and os.path.exists(template_path):
        prs = Presentation(template_path)
        # 기존 슬라이드 제거
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].get(
                '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
            )
            if rId:
                prs.part.drop_rel(rId)
            prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])
        print(f"  템플릿: {template_path}")
    else:
        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT
        print("  새 프레젠테이션 (템플릿 없음)")

    # 3. 슬라이드 빌드
    print("\n[3/4] 슬라이드 생성 중...")

    slide_count = 0

    # === 표지 3장 (본문 슬라이드와 동일한 네이티브 디자인으로 통일) ===
    # 2026-06-06 신교수님 지시: 표지도 검정 배경+좌측 황금 바+진네이비 밴드+흰 나눔명조로 통일.
    # 표지 1: 성경 말씀 (밴드="성경 말씀" / 본문=책명+장절, 책명 노랑 강조)
    build_native_text_slide(prs, "성경 말씀", f"{book_name}\n{chapter_verse}", [book_name])
    slide_count += 1
    print(f"  S{slide_count:02d} native_cover(말씀)  {book_name} {chapter_verse}")

    # 표지 2: 성경 본문 (원고에서 추출 시도)
    scripture_text = find_scripture_text(sermon_dir)
    if scripture_text:
        # 본문 끝의 성경구절 표기 제거 — 상단 밴드에 이미 표시됨 (신교수님 지시 2026-06-06)
        scripture_text = re.sub(r'\s*\([^()]*\d+\s*:\s*[\d,\s\-]+\)\s*$', '', scripture_text).rstrip()
        cover2_title = ref_label or f"{book_name} {chapter_verse}"
        build_native_text_slide(prs, cover2_title, scripture_text, [])
        slide_count += 1
        print(f"  S{slide_count:02d} native_cover(본문)  {cover2_title}")

    # 표지 3: 말씀 제목 (밴드="말씀 제목" / 본문=제목, 노랑 강조)
    build_native_text_slide(prs, "말씀 제목", title, [title])
    slide_count += 1
    print(f"  S{slide_count:02d} native_cover(제목)  {title}")

    # === 본문 슬롯 ===
    for slot in content_slots:
        num = slot.get('num', '??')
        desc = slot.get('desc', '')
        has_image = slot.get('has_image', False)
        text_info = slot.get('text', {}) if isinstance(slot.get('text'), dict) else {}
        masks = _extract_masks(slot)

        if has_image:
            img_path = find_image(sermon_dir, num, True)
            if img_path:
                if masks:
                    mask_positions = [
                        (
                            mask.get('type', 'rect'),
                            mask.get('left_cm', 0),
                            mask.get('top_cm', 0),
                            mask.get('width_cm', 0),
                            mask.get('height_cm', 0),
                        )
                        for mask in masks
                    ]
                    build_image_with_masks(prs, img_path, mask_positions)
                    slide_kind = "image_with_masks"
                else:
                    build_full_image(prs, img_path)
                    slide_kind = "full_image"
                slide_count += 1
                print(f"  S{slide_count:02d} {slide_kind:<18} PP{num} — {desc[:30]}")
            else:
                print(f"  ⚠ PP{num} 이미지 없음: kids_ppt_{num}.png")
        else:
            # 본문 슬롯 = 수정 가능한 네이티브 텍스트 (2026-06-06 신교수님 지시).
            # 디자인·글씨체·강조는 미리보기 본문 슬라이드(text_PP.png)와 픽셀 일치하되
            # PowerPoint에서 글자를 직접 수정할 수 있다. text 없으면 기존 이미지로 폴백.
            slot_title = text_info.get('title', '') if text_info else ''
            slot_body = text_info.get('body', '') if text_info else ''
            hl = _extract_highlights(text_info)
            if slot_title or slot_body:
                build_native_text_slide(prs, slot_title, slot_body, hl)
                slide_count += 1
                print(f"  S{slide_count:02d} native_text        PP{num} — {desc[:30]}")
            else:
                img_path = find_image(sermon_dir, num, False)
                if img_path:
                    build_full_image(prs, img_path)
                    slide_count += 1
                    print(f"  S{slide_count:02d} full_image(text)   PP{num} — {desc[:30]}")
                else:
                    print(f"  ⚠ PP{num} text/이미지 없음: 스킵")

    # 4. 저장
    print("\n[4/4] 저장 중...")

    # 파일명 생성
    folder_name = sermon_dir.name
    # 파일명 안전화: Windows 금지문자 제거, 비면 폴더명 제목 사용
    _safe = re.sub(r'[\\/:*?"<>|]', '', str(title)).strip()
    if not _safe:
        _safe = sermon_dir.name.split('-')[-1]
    # 부서 자동 감지 (폴더명 기준) — 아동부면 "아동부-" 접두, 장년부 등은 제목만
    _dept = next((d for d in ['아동부', '청소년부', '새벽', '금요'] if d in sermon_dir.name), '')
    output_name = f"{_dept + '-' if _dept else ''}{_safe}.pptx"
    output_path = output_dir / output_name

    prs.save(str(output_path))

    print(f"\n{'=' * 60}")
    print(f"아동부 PPT 생성 완료!")
    print(f"파일: {output_path}")
    print(f"슬라이드: {slide_count}장")
    print(f"{'=' * 60}")

    return str(output_path)


if __name__ == '__main__':
    if len(sys.argv) < 2:
        print("사용법: python scripts/sermon_pptx.py <설교폴더>")
        print("  예: python scripts/sermon_pptx.py sermons/20260607-고전13장-우산-아동부")
        sys.exit(1)

    sermon_dir = sys.argv[1]

    # 절대 경로 변환
    if not os.path.isabs(sermon_dir):
        sermon_dir = os.path.join(os.getcwd(), sermon_dir)

    template = sys.argv[2] if len(sys.argv) > 2 else None

    generate_kids_pptx(sermon_dir, template)
